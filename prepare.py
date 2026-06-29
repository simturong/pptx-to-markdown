"""
PPTX 파싱 스크립트
- shape 좌표(top, left) 기준 위→아래, 좌→우 순서로 요소 정렬
- 텍스트 추출 (bold/italic 보존)
- 이미지 추출 → output/images/ (v2: 픽셀 크기 포함)
- output/manifest.json 생성 (전체 합본 1개)

사용법:
  python prepare.py --check              # 의존성 확인
  python prepare.py <파일.pptx> [출력폴더]
"""

import json
import sys
from pathlib import Path


def check_dependencies():
    """필수 라이브러리 설치 여부 확인."""
    required = {"python-pptx": "pptx", "Pillow": "PIL"}
    missing = []
    for pkg, module in required.items():
        try:
            __import__(module)
            print(f"  [OK] {pkg}")
        except ImportError:
            print(f"  [없음] {pkg}")
            missing.append(pkg)
    if missing:
        print(f"\n설치 명령어:")
        print(f'  C:\\Users\\sando\\AppData\\Local\\Programs\\Python\\Python313\\python.exe -m pip install {" ".join(missing)}')
        sys.exit(1)
    else:
        print("\n모든 의존성 확인 완료.")


if len(sys.argv) == 2 and sys.argv[1] == "--check":
    print("의존성 확인 중...")
    check_dependencies()
    sys.exit(0)

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("[오류] python-pptx 가 설치되어 있지 않습니다.")
    print("실행: C:\\Users\\sando\\AppData\\Local\\Programs\\Python\\Python313\\python.exe -m pip install python-pptx Pillow")
    sys.exit(1)

try:
    from PIL import Image as PILImage
    _PILLOW_AVAILABLE = True
except ImportError:
    _PILLOW_AVAILABLE = False


def emu_to_pt(emu):
    return round(emu / 12700, 1)


def extract_text(shape):
    """텍스트 shape에서 내용 추출. bold/italic 마크다운 변환."""
    lines = []
    for para in shape.text_frame.paragraphs:
        parts = []
        for run in para.runs:
            text = run.text
            if not text.strip():
                continue
            if run.font.bold and run.font.italic:
                text = f"***{text}***"
            elif run.font.bold:
                text = f"**{text}**"
            elif run.font.italic:
                text = f"*{text}*"
            parts.append(text)
        line = "".join(parts)
        if line:
            lines.append(line)
    return "\n".join(lines)


def is_title(shape):
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        return shape.placeholder_format and shape.placeholder_format.idx == 0
    except Exception:
        return False


def extract_image(shape, slide_idx, img_idx, image_dir):
    """이미지 shape에서 PNG 파일 저장. (파일명, 너비px, 높이px) 반환."""
    try:
        image = shape.image
        ext = image.ext  # png, jpg, gif, ...
        filename = f"{slide_idx:02d}_{img_idx:02d}.{ext}"
        path = image_dir / filename
        path.write_bytes(image.blob)

        width_px, height_px = None, None
        if _PILLOW_AVAILABLE:
            try:
                with PILImage.open(path) as img:
                    width_px, height_px = img.size
            except Exception:
                pass

        return filename, width_px, height_px
    except Exception as e:
        print(f"  [경고] 이미지 추출 실패: {e}", file=sys.stderr)
        return None, None, None


def sort_key(shape):
    """위→아래, 좌→우 정렬 키."""
    top = shape.top if shape.top is not None else 0
    left = shape.left if shape.left is not None else 0
    # 슬라이드 높이 기준 행 구분 (±5% 오차 허용)
    row_bucket = round(top / 457200)  # 457200 EMU ≈ 0.5인치 단위로 버킷
    return (row_bucket, left)


def process_slide(slide, slide_idx, image_dir):
    """슬라이드 하나를 파싱해서 elements 리스트 반환."""
    elements = []
    img_idx = 1

    shapes = sorted(slide.shapes, key=sort_key)

    for shape in shapes:
        pos = {
            "top_pt": emu_to_pt(shape.top or 0),
            "left_pt": emu_to_pt(shape.left or 0),
        }

        # 텍스트 shape
        if shape.has_text_frame:
            text = extract_text(shape)
            if not text.strip():
                continue
            kind = "title" if is_title(shape) else "text"
            elements.append({
                "type": kind,
                "text": text,
                "position": pos,
            })

        # 이미지 shape
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            filename, width_px, height_px = extract_image(shape, slide_idx, img_idx, image_dir)
            if filename:
                img_entry = {
                    "type": "image",
                    "path": f"images/{filename}",
                    "position": pos,
                }
                if width_px is not None:
                    img_entry["width_px"] = width_px
                    img_entry["height_px"] = height_px
                if shape.width and shape.height:
                    img_entry["display_width_pt"] = emu_to_pt(shape.width)
                    img_entry["display_height_pt"] = emu_to_pt(shape.height)
                elements.append(img_entry)
                img_idx += 1

        # 그룹 shape (재귀 처리 생략, 필요 시 확장)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            elements.append({
                "type": "note",
                "text": "[그룹 shape — 수동 확인 필요]",
                "position": pos,
            })

    return elements


def prepare(pptx_path: str, output_dir: str = "output"):
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    image_dir = output_dir / "images"

    image_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(pptx_path)
    total = len(prs.slides)
    print(f"총 슬라이드 수: {total}")

    all_slides = []

    for idx, slide in enumerate(prs.slides, start=1):
        print(f"  슬라이드 {idx}/{total} 처리 중...")
        elements = process_slide(slide, idx, image_dir)

        slide_data = {
            "slide": idx,
            "elements": elements,
        }
        all_slides.append(slide_data)

    # 전체 manifest
    full_manifest = output_dir / "manifest.json"
    full_manifest.write_text(
        json.dumps(all_slides, ensure_ascii=False, indent=2),
        encoding="utf-8"
    )

    print(f"\n완료!")
    print(f"  이미지   → {image_dir}")
    print(f"  manifest → {full_manifest}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("사용법: python prepare.py <파일.pptx> [출력폴더]")
        sys.exit(1)

    pptx_file = sys.argv[1]
    out_dir = sys.argv[2] if len(sys.argv) > 2 else "output"
    prepare(pptx_file, out_dir)
