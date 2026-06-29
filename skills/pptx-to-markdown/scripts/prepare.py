"""
PPTX Parsing Script
- Sort elements based on shape coordinates (top, left) in top-to-bottom, left-to-right order.
- Extract text (preserving bold/italic markdown formatting).
- Extract images → output/images/ (includes pixel size in v2).
- Generate output/manifest.json (single unified file).

Usage:
  python prepare.py --check              # Check dependencies
  python prepare.py <file.pptx> [output_dir]
"""

import json
import sys
from pathlib import Path


def check_dependencies():
    """Verify if required packages are installed."""
    required = {"python-pptx": "pptx", "Pillow": "PIL"}
    missing = []
    for pkg, module in required.items():
        try:
            __import__(module)
            print(f"  [OK] {pkg}")
        except ImportError:
            print(f"  [Missing] {pkg}")
            missing.append(pkg)
    if missing:
        print(f"\nInstall command:")
        print(f'  C:\\Users\\sando\\AppData\\Local\\Programs\\Python\\Python313\\python.exe -m pip install {" ".join(missing)}')
        sys.exit(1)
    else:
        print("\nAll dependencies are verified.")


if len(sys.argv) == 2 and sys.argv[1] == "--check":
    print("Checking dependencies...")
    check_dependencies()
    sys.exit(0)

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("[Error] python-pptx is not installed.")
    print("Run: C:\\Users\\sando\\AppData\\Local\\Programs\\Python\\Python313\\python.exe -m pip install python-pptx Pillow")
    sys.exit(1)

try:
    from PIL import Image as PILImage
    _PILLOW_AVAILABLE = True
except ImportError:
    _PILLOW_AVAILABLE = False


def emu_to_pt(emu):
    return round(emu / 12700, 1)


def extract_text(shape):
    """Extract text from text shape and apply bold/italic markdown formatting."""
    lines = []
    for para in shape.text_frame.paragraphs:
        parts = []
        for run in para.runs:
            text = run.text
            if not text.strip():
                continue
            if run.font.bold and run.font.italic:
                text = f"***{text}***"
            elif run.font.bold:
                text = f"**{text}**"
            elif run.font.italic:
                text = f"*{text}*"
            parts.append(text)
        line = "".join(parts)
        if line:
            lines.append(line)
    return "\n".join(lines)


def is_title(shape):
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        return shape.placeholder_format and shape.placeholder_format.idx == 0
    except Exception:
        return False


def extract_image(shape, slide_idx, img_idx, image_dir):
    """Save image shape to PNG file. Returns (filename, width_px, height_px)."""
    try:
        image = shape.image
        ext = image.ext  # png, jpg, gif, ...
        filename = f"{slide_idx:02d}_{img_idx:02d}.{ext}"
        path = image_dir / filename
        path.write_bytes(image.blob)

        width_px, height_px = None, None
        if _PILLOW_AVAILABLE:
            try:
                with PILImage.open(path) as img:
                    width_px, height_px = img.size
            except Exception:
                pass

        return filename, width_px, height_px
    except Exception as e:
        print(f"  [Warning] Failed to extract image: {e}", file=sys.stderr)
        return None, None, None


def sort_key(shape):
    """Sorting key for top-to-bottom, left-to-right order."""
    top = shape.top if shape.top is not None else 0
    left = shape.left if shape.left is not None else 0
    # Group shapes into row buckets of 0.5 inches (457200 EMU) to tolerate minor alignment issues
    row_bucket = round(top / 457200)
    return (row_bucket, left)


def process_slide(slide, slide_idx, image_dir):
    """Parse a single slide and return its elements list."""
    elements = []
    img_idx = 1

    shapes = sorted(slide.shapes, key=sort_key)

    for shape in shapes:
        pos = {
            "top_pt": emu_to_pt(shape.top or 0),
            "left_pt": emu_to_pt(shape.left or 0),
        }

        # Text shape
        if shape.has_text_frame:
            text = extract_text(shape)
            if not text.strip():
                continue
            kind = "title" if is_title(shape) else "text"
            elements.append({
                "type": kind,
                "text": text,
                "position": pos,
            })

        # Image shape
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            filename, width_px, height_px = extract_image(shape, slide_idx, img_idx, image_dir)
            if filename:
                img_entry = {
                    "type": "image",
                    "path": f"images/{filename}",
                    "position": pos,
                }
                if width_px is not None:
                    img_entry["width_px"] = width_px
                    img_entry["height_px"] = height_px
                if shape.width and shape.height:
                    img_entry["display_width_pt"] = emu_to_pt(shape.width)
                    img_entry["display_height_pt"] = emu_to_pt(shape.height)
                elements.append(img_entry)
                img_idx += 1

        # Group shape
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            elements.append({
                "type": "note",
                "text": "[Group shape — manual verification required]",
                "position": pos,
            })

    return elements


def prepare(pptx_path: str, output_dir: str = "output"):
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    image_dir = output_dir / "images"

    image_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(pptx_path)
    total = len(prs.slides)
    print(f"Total slides: {total}")

    all_slides = []

    for idx, slide in enumerate(prs.slides, start=1):
        print(f"  Processing slide {idx}/{total}...")
        elements = process_slide(slide, idx, image_dir)

        slide_data = {
            "slide": idx,
            "elements": elements,
        }
        all_slides.append(slide_data)

    # Save manifest
    full_manifest = output_dir / "manifest.json"
    full_manifest.write_text(
        json.dumps(all_slides, ensure_ascii=False, indent=2),
        encoding="utf-8"
    )

    print(f"\nCompleted!")
    print(f"  Images   → {image_dir}")
    print(f"  Manifest → {full_manifest}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python prepare.py <file.pptx> [output_dir]")
        sys.exit(1)

    pptx_file = sys.argv[1]
    out_dir = sys.argv[2] if len(sys.argv) > 2 else "output"
    prepare(pptx_file, out_dir)
